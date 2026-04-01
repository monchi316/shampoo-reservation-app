#!/usr/bin/env python3
"""社内運用向け（管理者マニュアル寄り）PowerPoint を生成する。"""
from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "社内運用_管理者マニュアル.pptx"

JP_FONT = "Yu Gothic"  # 環境により置換される場合あり


def _set_run_font(run, size_pt: int = 18, bold: bool = False) -> None:
    run.font.name = JP_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            _set_run_font(r, 32, True)
    for p in sub.text_frame.paragraphs:
        for r in p.runs:
            _set_run_font(r, 16)


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            _set_run_font(r, 28, True)
    body = slide.placeholders[1].text_frame
    body.clear()
    first = True
    for line in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        for r in p.runs:
            _set_run_font(r, 18)


def add_two_column_note_slide(prs: Presentation, title: str, left: List[str], right: List[str]) -> None:
    """図の代わりに左右2列で整理（社内マニュアル向け）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = tx.text_frame
    tf.text = title
    for p in tf.paragraphs:
        for r in p.runs:
            _set_run_font(r, 28, True)

    lx = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5))
    ltf = lx.text_frame
    ltf.clear()
    first = True
    for line in left:
        p = ltf.paragraphs[0] if first else ltf.add_paragraph()
        first = False
        p.text = line
        for r in p.runs:
            _set_run_font(r, 16)
    rx = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.5))
    rtf = rx.text_frame
    rtf.clear()
    first = True
    for line in right:
        p = rtf.paragraphs[0] if first else rtf.add_paragraph()
        first = False
        p.text = line
        for r in p.runs:
            _set_run_font(r, 16)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "洗車予約システム（社内運用）",
        "管理者マニュアル（運用・設定・障害対応のたたき台）",
    )

    add_bullet_slide(
        prs,
        "この資料の目的",
        [
            "社内の運用担当者が、日々の確認と設定変更を迷わず行えるようにする",
            "「どこを触ると何が変わるか」を明確にする",
            "障害時に「まず何を見るか」を揃える",
            "※ コードの書き方は扱わない（画面と設定中心）",
        ],
    )

    add_bullet_slide(
        prs,
        "対象者・前提",
        [
            "対象：社内の運用担当、サポート担当、導入担当",
            "前提：本番URL・管理画面URL・Supabase・LINEの管理権限がある",
            "現在の構成：LINE（LIFF）＋ Vercel（Web）＋ Supabase（DB/Storage）",
        ],
    )

    add_two_column_note_slide(
        prs,
        "全体像（役割のイメージ）",
        [
            "【LINE】",
            "・ユーザー入口（予約画面を開く）",
            "・通知（予約/変更/キャンセル）",
            "",
            "【Vercel】",
            "・予約アプリ本体を公開",
            "・サーバーAPI（通知・予約書き込み等）",
        ],
        [
            "【Supabase】",
            "・予約データ保存",
            "・店舗設定（メニュー/営業時間）",
            "・ロゴ画像（Storage）",
            "",
            "【tenant_id】",
            "・企業（店舗）単位でデータを分けるID",
            "・現状は既定テナントを利用（将来はログイン連携で切替）",
        ],
    )

    add_bullet_slide(
        prs,
        "管理者が触る主な画面",
        [
            "管理トップ：/admin",
            "予約一覧：/admin/reservations",
            "予約詳細：/admin/reservations/[id]",
            "売上：/admin/sales",
            "店舗設定：/admin/settings（ロゴ・メニュー・営業時間・予約枠）",
        ],
    )

    add_bullet_slide(
        prs,
        "日次チェック（推奨）",
        [
            "予約一覧で「当日・翌日」の件数をざっと確認",
            "キャンセル・変更が入っていないか（通知トラブルの早期発見）",
            "予約詳細でステータスが想定どおりか（未対応が滞留していないか）",
        ],
    )

    add_bullet_slide(
        prs,
        "週次チェック（推奨）",
        [
            "売上画面で完了件数の傾向を確認",
            "店舗設定の例外日（臨時休業）が翌週分まで入っているか",
            "メニュー料金の変更予定があれば反映済みか",
        ],
    )

    add_bullet_slide(
        prs,
        "月次チェック（推奨）",
        [
            "営業時間・平均作業時間・移動時間の見直し（混雑/遅延の実態ベース）",
            "ロゴ差し替え（キャンペーン・ブランド刷新）",
            "LINE側の導線（リッチメニュー）文言・リンクの棚卸し",
        ],
    )

    add_bullet_slide(
        prs,
        "店舗設定：ロゴ",
        [
            "場所：/admin/settings → 店舗ロゴ",
            "対応形式：png / jpg / webp（目安5MB以下）",
            "反映先：予約ページ左上（枠内に自動収まる表示）",
            "トラブル：壊れた画像 → Storage/URL周りを疑う（再アップロード）",
        ],
    )

    add_bullet_slide(
        prs,
        "店舗設定：メニュー・料金",
        [
            "料金の正はDB（service_menu_items）",
            "S/M/L・内装は slug（size_s 等）で予約画面と連動",
            "追加メニューは行を足せるが、予約合計に載せるには画面側対応が別途必要",
            "保存後、予約画面をリロードして表示確認",
        ],
    )

    add_bullet_slide(
        prs,
        "店舗設定：営業時間（基本）",
        [
            "モード：全日同一 / 曜日別",
            "例外日：臨時休業・その日だけの営業時間",
            "保存忘れに注意（保存ボタンは営業・予約設定セクション）",
        ],
    )

    add_bullet_slide(
        prs,
        "店舗設定：予約枠の意味（重要）",
        [
            "1台あたり平均作業時間 × 台数 = 占有時間",
            "平均移動時間 = 前後の予約との間に空けるバッファ",
            "「営業時間内」でも、終了が閉店を超えると不可になる",
            "変更時はユーザーへ周知（特に作業時間を伸ばした場合）",
        ],
    )

    add_bullet_slide(
        prs,
        "予約一覧の見方",
        [
            "日付・ステータスで絞り込み",
            "同一グループ（複数台）は group_id でまとまる想定",
            "詳細へ進んで金額提案・メモ・担当者名などを更新",
        ],
    )

    add_bullet_slide(
        prs,
        "予約詳細でよくやる操作",
        [
            "ステータス更新（確認済/完了/キャンセル等）",
            "売上金額の入力（複数台は行ごと）",
            "メモ・スタッフ名・作業完了日時の記録",
        ],
    )

    add_bullet_slide(
        prs,
        "売上画面",
        [
            "完了ステータスを集計して参照する想定",
            "数字が合わない場合は「ステータス定義」と「対象期間」を先に確認",
        ],
    )

    add_bullet_slide(
        prs,
        "ユーザー向けURL（社内周知テンプレ）",
        [
            "予約：本番の /reserve（LIFFから開くのが基本）",
            "変更・キャンセル：LINE通知に付くURL（id / groupId 付き）",
            "※ URLを社外に転送されると閲覧リスクがある点を理解して運用",
        ],
    )

    add_bullet_slide(
        prs,
        "tenant_id / 環境変数（社内メモ）",
        [
            "NEXT_PUBLIC_DEFAULT_TENANT_ID：既定の店舗ID（本番でズレるとデータが見えない）",
            "Supabase の service_role はサーバーのみ（漏洩厳禁）",
            "LINE_CHANNEL_ACCESS_TOKEN：通知が止まるとまず疑う",
        ],
    )

    add_bullet_slide(
        prs,
        "障害対応：予約できない",
        [
            "まず予約画面でアラート文言を確認（営業外/重複/設定なし）",
            "店舗設定の営業時間・例外日・作業/移動時間を確認",
            "API 500 の場合は Vercel ログと Supabase 側エラーをエスカレーション",
        ],
    )

    add_bullet_slide(
        prs,
        "障害対応：LINE通知が来ない",
        [
            "LINE_CHANNEL_ACCESS_TOKEN（本番環境変数）",
            "ユーザーID（LINE user id）が取れているか",
            "キャンセル/予約APIのレスポンスが成功か",
        ],
    )

    add_bullet_slide(
        prs,
        "障害対応：画像（ロゴ）が壊れる",
        [
            "再アップロード",
            "Storage bucket（tenant-logos）とDBの logo_path",
            "署名付きURLの期限切れを疑う場合は再読み込みで更新されるか確認",
        ],
    )

    add_bullet_slide(
        prs,
        "障害対応：DB/権限（RLS）系",
        [
            "ブラウザ直叩きで弾かれる場合はRLSやAPI経由の設計差を疑う",
            "エラーコード 42501 は権限（RLS）が典型",
            "最終的にシステム担当へ（ポリシー/ルートの見直し）",
        ],
    )

    add_two_column_note_slide(
        prs,
        "エスカレーションフロー（例）",
        [
            "【レベル1：運用】",
            "・設定誤りの修正",
            "・再現手順のメモ",
            "・スクリーンショット取得",
        ],
        [
            "【レベル2：システム】",
            "・環境変数/デプロイ",
            "・Supabase SQL/Storage",
            "・LINE Developers設定",
            "",
            "【連絡時に添える情報】",
            "時刻、URL、ユーザーID、予約ID、画面録画",
        ],
    )

    add_bullet_slide(
        prs,
        "セキュリティ運用上の注意",
        [
            "管理画面は現状ガードが弱い前提でURLの取り扱いに注意",
            "service_role キーをチャットに貼らない",
            "権限最小化（社内でも閲覧者/編集者を分ける）",
        ],
    )

    add_bullet_slide(
        prs,
        "改定・配布のしかた",
        [
            "改定履歴（日付・変更者・要点）を表紙裏または最終スライドに残す",
            "社内ポータルに最新版を1つに集約",
            "対外向け資料とは分ける（この資料は社内運用寄り）",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
